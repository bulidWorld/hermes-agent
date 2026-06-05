"""
File type classifier and best-effort document text extractor.

Pure logic — no network, no state, no mandatory imports beyond the
standard library.  Parsing libraries (``python-docx``, ``pypdf``,
``openpyxl``, ``python-pptx``) are imported lazily so they remain
optional.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional, FrozenSet

logger = logging.getLogger(__name__)


class FileClassifier:
    """Classify files by MIME type / extension and extract text from documents.

    All methods are ``@staticmethod`` — the class is a namespace, not a
    stateful object.
    """

    TEXT_INLINE_LIMIT: int = 100_000  # 100 KB

    # ---- text MIME types ---------------------------------------------------
    TEXT_MIME: FrozenSet[str] = frozenset({
        "application/json",
        "application/xml",
        "application/javascript",
        "application/x-yaml",
        "application/x-httpd-php",
    })

    # ---- document extensions & MIME types that can be parsed ---------------
    DOC_EXT: FrozenSet[str] = frozenset({".docx", ".pdf", ".xlsx", ".pptx"})
    DOC_MIME: FrozenSet[str] = frozenset({
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/msword",
        "application/vnd.ms-excel",
        "application/vnd.ms-powerpoint",
    })

    # ---- code file extensions (treated as text) ----------------------------
    CODE_EXT: FrozenSet[str] = frozenset({
        ".py", ".js", ".ts", ".go", ".rs", ".java", ".c", ".cpp", ".h",
        ".sh", ".bash", ".zsh", ".yaml", ".yml", ".toml", ".ini", ".cfg",
        ".conf", ".md", ".rst", ".sql", ".rb", ".php", ".swift", ".kt",
        ".scala", ".lua", ".r",
    })

    # ------------------------------------------------------------------
    # Classification
    # ------------------------------------------------------------------

    @staticmethod
    def classify(mime_type: str, filename: str) -> str:
        """Return one of ``"text"``, ``"image"``, ``"document"``, or
        ``"binary"``.

        Priority: image > document > text > binary.
        """
        mime = (mime_type or "").strip().lower()
        ext = Path(filename or "").suffix.lower()

        # 1. Images
        if mime.startswith("image/"):
            logger.warning("classify → image (mime=%r)", mime)
            return "image"

        # 2. Parseable documents
        if ext in FileClassifier.DOC_EXT or mime in FileClassifier.DOC_MIME:
            logger.warning(
                "classify → document (ext=%r in DOC_EXT=%s, mime=%r in DOC_MIME=%s)",
                ext, ext in FileClassifier.DOC_EXT, mime, mime in FileClassifier.DOC_MIME,
            )
            return "document"

        # 3. Known text types
        if mime.startswith("text/") or mime in FileClassifier.TEXT_MIME:
            logger.warning("classify → text (mime=%r)", mime)
            return "text"

        # 4. Code files by extension
        if ext in FileClassifier.CODE_EXT:
            logger.warning("classify → text via code ext (ext=%r)", ext)
            return "text"

        # 5. Everything else
        logger.warning("classify → binary (mime=%r, ext=%r)", mime, ext)
        return "binary"

    # ------------------------------------------------------------------
    # Document text extraction
    # ------------------------------------------------------------------

    # MIME → parser method name (fallback when the cached file path has no
    # extension — the cache key is a generated id, not the original filename).
    _MIME_TO_PARSER: dict[str, str] = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "_parse_docx",
        "application/msword": "_parse_docx",
        "application/pdf": "_parse_pdf",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "_parse_xlsx",
        "application/vnd.ms-excel": "_parse_xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "_parse_pptx",
        "application/vnd.ms-powerpoint": "_parse_pptx",
    }

    @staticmethod
    def parse_document(file_path: str, mime_type: str) -> Optional[str]:
        """Best-effort text extraction from common document formats.

        ====== ===========================================
        Format Library (lazy-imported)
        ====== ===========================================
        .docx  ``python-docx``
        .pdf   ``pypdf`` (fallback: ``PyPDF2``)
        .xlsx  ``openpyxl``
        .pptx  ``python-pptx``
        ====== ===========================================

        Routing is first by file extension, then by MIME type as a
        fallback (cached files may be stored without their original
        extension).  Returns ``None`` when the required library is not
        installed or the file is unreadable — callers fall back to
        path-only injection.
        """
        ext = Path(file_path).suffix.lower()

        # -- extension-based routing ----------------------------------------
        if ext == ".docx":
            logger.warning("parse_document → .docx path (%s)", file_path)
            return FileClassifier._parse_docx(file_path)
        if ext == ".pdf":
            logger.warning("parse_document → .pdf path (%s)", file_path)
            return FileClassifier._parse_pdf(file_path)
        if ext == ".xlsx":
            logger.warning("parse_document → .xlsx path (%s)", file_path)
            return FileClassifier._parse_xlsx(file_path)
        if ext == ".pptx":
            logger.warning("parse_document → .pptx path (%s)", file_path)
            return FileClassifier._parse_pptx(file_path)

        # -- MIME-based fallback (cached files may lack extension) ----------
        parser_name = FileClassifier._MIME_TO_PARSER.get(mime_type)
        if parser_name is not None:
            logger.warning(
                "parse_document → ext=%r missing, routing by mime=%r → %s (%s)",
                ext, mime_type, parser_name, file_path,
            )
            parser = getattr(FileClassifier, parser_name)
            return parser(file_path)

        logger.warning(
            "parse_document → unsupported ext=%r mime=%r, returning None (%s)",
            ext, mime_type, file_path,
        )
        return None

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def format_size(size_bytes: int) -> str:
        """Human-readable size string, e.g. ``"2.4MB"``."""
        if size_bytes >= 1_000_000:
            return f"{size_bytes / 1_000_000:.1f}MB"
        if size_bytes >= 1_000:
            return f"{size_bytes / 1_000:.0f}KB"
        return f"{size_bytes}B"

    # ==================================================================
    # Internal parsers (lazy imports, all-errors-caught)
    # ==================================================================

    @staticmethod
    def _parse_docx(file_path: str) -> Optional[str]:
        try:
            from docx import Document  # type: ignore[import-untyped]
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            result = "\n".join(paragraphs) if paragraphs else None
            logger.warning("_parse_docx → %s (%d paragraphs)", "success" if result else "empty", len(paragraphs))
            return result
        except ImportError:
            logger.warning("python-docx not installed — cannot parse .docx")
            return None
        except Exception as exc:
            logger.warning("Failed to parse .docx %s: %s", file_path, exc)
            return None

    @staticmethod
    def _parse_pdf(file_path: str) -> Optional[str]:
        # Prefer pypdf, fall back to PyPDF2
        for lib_name, import_path in (
            ("pypdf", "pypdf"),
            ("PyPDF2", "PyPDF2"),
        ):
            try:
                module = __import__(import_path, fromlist=["PdfReader"])
                reader = module.PdfReader(file_path)
                logger.warning("_parse_pdf → using %s for %s", lib_name, file_path)
                pages = [
                    (page.extract_text() or "")
                    for page in reader.pages
                ]
                text = "\n\n".join(p for p in pages if p.strip())
                result = text if text.strip() else None
                logger.warning("_parse_pdf → %s (%d pages)", "success" if result else "empty", len(pages))
                return result
            except ImportError:
                logger.warning("%s not installed, trying next fallback", lib_name)
                continue
            except Exception as exc:
                logger.warning("%s failed for %s: %s", lib_name, file_path, exc)
                continue
        logger.warning("No PDF library available (pypdf / PyPDF2) — cannot parse .pdf")
        return None

    @staticmethod
    def _parse_xlsx(file_path: str) -> Optional[str]:
        try:
            from openpyxl import load_workbook  # type: ignore[import-untyped]
            wb = load_workbook(file_path, read_only=True, data_only=True)
            parts: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"[Sheet: {sheet_name}]")
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        parts.append("\t".join(cells))
                        row_count += 1
                    if row_count >= 5000:
                        parts.append("[Sheet truncated at 5000 rows]")
                        break
            wb.close()
            result = "\n".join(parts) if parts else None
            logger.warning("_parse_xlsx → %s (%d sheets)", "success" if result else "empty", len(wb.sheetnames))
            return result
        except ImportError:
            logger.warning("openpyxl not installed — cannot parse .xlsx")
            return None
        except Exception as exc:
            logger.warning("Failed to parse .xlsx %s: %s", file_path, exc)
            return None

    @staticmethod
    def _parse_pptx(file_path: str) -> Optional[str]:
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
            prs = Presentation(file_path)
            parts: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            t = paragraph.text.strip()
                            if t:
                                texts.append(t)
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            result = "\n\n".join(parts) if parts else None
            logger.warning("_parse_pptx → %s (%d slides with text)", "success" if result else "empty", len(parts))
            return result
        except ImportError:
            logger.warning("python-pptx not installed — cannot parse .pptx")
            return None
        except Exception as exc:
            logger.warning("Failed to parse .pptx %s: %s", file_path, exc)
            return None
